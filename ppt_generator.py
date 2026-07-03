"""
PPT_GENERATOR.PY V5.1
BASE: V5.0 + Integración sistema de caché

NUEVO V5.1:
- Import de gif_cache (los GIFs ya están cacheados por timelapse_generator.py)
- Mensaje informativo sobre GIFs encontrados
- NO genera GIFs (solo los usa)
"""

import os
import glob
import copy
import shutil
import zipfile
import tempfile
from lxml import etree
from datetime import datetime
from pptx import Presentation
from PIL import Image

# =========================
# NUEVO V5.1: IMPORT SISTEMA DE CACHÉ
# =========================
# Nota: ppt_generator NO genera GIFs, solo los usa
# El caché ya fue implementado en timelapse_generator.py
from gif_cache import existe_en_cache

VOLCANES_ACTIVOS = [
    # ZONA NORTE
    "Taapaca", "Parinacota", "Guallatiri", "Isluga", "Irruputuncu", "Ollague", "San Pedro", "Lascar",
    # ZONA CENTRO
    "Tupungatito", "San Jose", "Tinguiririca", "Planchon-Peteroa", "Descabezado Grande", 
    "Tatara-San Pedro", "Laguna del Maule", "Nevado de Longavi", "Nevados de Chillan",
    # ZONA SUR
    "Antuco", "Copahue", "Callaqui", "Lonquimay", "Llaima", "Sollipulli", "Villarrica", 
    "Quetrupillan", "Lanin", "Mocho-Choshuenco", "Carran - Los Venados", "Puyehue - Cordon Caulle", 
    "Antillanca - Casablanca",
    # ZONA AUSTRAL
    "Osorno", "Calbuco", "Yate", "Hornopiren", "Huequi", "Michinmahuida", "Chaiten",
    "Corcovado", "Melimoyu", "Mentolat", "Cay", "Maca", "Hudson",
    # VISTAS ZOOM ADICIONALES (SERNAGEOMIN)
    "Melimoyu_Conos_Eruptivos", "Mentolat_Sismicidad_VT", "Hudson_Ultima_Erupcion"
]
PLANTILLA_PATH = "docs/plantillas/Cambios_morfologicos.pptx"
OUTPUT_DIR = "docs/sentinel2"

MESES_ES = {
    1: 'enero', 2: 'febrero', 3: 'marzo', 4: 'abril',
    5: 'mayo', 6: 'junio', 7: 'julio', 8: 'agosto',
    9: 'septiembre', 10: 'octubre', 11: 'noviembre', 12: 'diciembre'
}

def formatear_fecha_espanol(fecha_str):
    try:
        dt = datetime.strptime(fecha_str, '%Y-%m-%d')
        return f"{dt.strftime('%d')} {MESES_ES[dt.month]}"
    except:
        return fecha_str

def comprimir_gif(input_path, output_path, max_size_mb=1.0):
    try:
        size_mb = os.path.getsize(input_path) / (1024 * 1024)
        if size_mb <= max_size_mb:
            import shutil
            shutil.copy2(input_path, output_path)
            print(f"      ✅ GIF OK ({size_mb:.2f} MB)")
            return output_path
        
        print(f"      🔽 Comprimiendo ({size_mb:.2f} MB > {max_size_mb:.2f} MB)...")
        img = Image.open(input_path)
        frames = []
        try:
            while True:
                # Reducir colores a 64 (más agresivo)
                frame = img.copy().convert('P', palette=Image.ADAPTIVE, colors=64)
                frames.append(frame)
                img.seek(img.tell() + 1)
        except EOFError:
            pass
        
        # Reducir tamaño de frames si muy grande
        if size_mb > 2.0:
            new_size = (int(frames[0].width * 0.75), int(frames[0].height * 0.75))
            frames = [f.resize(new_size, Image.Resampling.LANCZOS) for f in frames]
        
        frames[0].save(output_path, save_all=True, append_images=frames[1:],
                      optimize=True, duration=img.info.get('duration', 100), loop=0)
        print(f"      ✅ {size_mb:.2f} MB → {os.path.getsize(output_path)/(1024*1024):.2f} MB")
        return output_path
    except Exception as e:
        print(f"      ❌ Error: {e}")
        import shutil
        shutil.copy2(input_path, output_path)
        return output_path

def generar_ppt(volcan_nombre):
    print(f"\n🌋 {volcan_nombre}")
    
    carpeta_timelapses = f"docs/sentinel2/{volcan_nombre}/timelapses_ppt"
    if not os.path.exists(carpeta_timelapses):
        print(f"    ❌ No existe: {carpeta_timelapses}")
        return None
    
    # =========================
    # NUEVO V5.1: BUSCAR GIFs (ya cacheados por timelapse_generator.py)
    # =========================
    gifs_rgb = sorted(glob.glob(f"{carpeta_timelapses}/{volcan_nombre}_RGB_*.gif"))
    gifs_thermal = sorted(glob.glob(f"{carpeta_timelapses}/{volcan_nombre}_ThermalFalseColor_*.gif"))
    
    print(f"    📂 GIFs encontrados: RGB={len(gifs_rgb)}, Thermal={len(gifs_thermal)}")
    
    if not gifs_rgb or not gifs_thermal:
        print(f"    ❌ GIFs incompletos")
        return None
    
    gif_rgb_path = gifs_rgb[-1]
    gif_thermal_path = gifs_thermal[-1]
    
    partes = os.path.basename(gif_rgb_path).replace('.gif', '').split('_')
    if len(partes) < 4:
        print(f"    ❌ No se pudieron extraer fechas")
        return None
    
    fecha_inicio, fecha_fin = partes[-2], partes[-1]
    print(f"    📅 {fecha_inicio} → {fecha_fin}")
    
    # tempdir portable (en Windows local no existe /tmp; en el runner Ubuntu es /tmp igual)
    temp_rgb = os.path.join(tempfile.gettempdir(), f"{volcan_nombre}_RGB.gif")
    temp_thermal = os.path.join(tempfile.gettempdir(), f"{volcan_nombre}_Thermal.gif")
    
    print(f"    🔽 Comprimiendo GIFs para PPT...")
    gif_rgb_final = comprimir_gif(gif_rgb_path, temp_rgb)
    gif_thermal_final = comprimir_gif(gif_thermal_path, temp_thermal)
    
    if not os.path.exists(PLANTILLA_PATH):
        print(f"    ❌ Plantilla no encontrada")
        return None
    
    prs = Presentation(PLANTILLA_PATH)
    slide = prs.slides[0]
    
    fecha_inicio_es = formatear_fecha_espanol(fecha_inicio)
    fecha_fin_es = formatear_fecha_espanol(fecha_fin)
    ano = fecha_fin.split('-')[0]
    
    texto_rgb = f"Imagenes Sentinel 2 L2A color verdadero, Time Lapse {fecha_inicio_es} → {fecha_fin_es} {ano}"
    texto_thermal = f"Imagenes Sentinel 2 L2A Falso color, Time Lapse {fecha_inicio_es} → {fecha_fin_es} {ano}"
    
    print(f"    ✏️  Actualizando textos...")
    textos_ok = 0
    
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not hasattr(shape, "text"):
            continue
        
        texto = shape.text.strip()
        
        # Reemplazar texto RGB timelapse
        if "color verdadero" in texto.lower() and "time lapse" in texto.lower():
            p = shape.text_frame.paragraphs[0]
            fmt = None
            if p.runs:
                fmt = {'name': p.runs[0].font.name, 'size': p.runs[0].font.size,
                      'bold': p.runs[0].font.bold, 'italic': p.runs[0].font.italic}
            p.clear()
            run = p.add_run()
            run.text = texto_rgb
            if fmt:
                if fmt['name']: run.font.name = fmt['name']
                if fmt['size']: run.font.size = fmt['size']
                if fmt['bold'] is not None: run.font.bold = fmt['bold']
                if fmt['italic'] is not None: run.font.italic = fmt['italic']
            print(f"       ✅ RGB")
            textos_ok += 1
        
        # Reemplazar texto Thermal timelapse
        elif "falso color" in texto.lower() and "time lapse" in texto.lower():
            p = shape.text_frame.paragraphs[0]
            fmt = None
            if p.runs:
                fmt = {'name': p.runs[0].font.name, 'size': p.runs[0].font.size,
                      'bold': p.runs[0].font.bold, 'italic': p.runs[0].font.italic}
            p.clear()
            run = p.add_run()
            run.text = texto_thermal
            if fmt:
                if fmt['name']: run.font.name = fmt['name']
                if fmt['size']: run.font.size = fmt['size']
                if fmt['bold'] is not None: run.font.bold = fmt['bold']
                if fmt['italic'] is not None: run.font.italic = fmt['italic']
            print(f"       ✅ Thermal")
            textos_ok += 1
        
        # Reemplazar nombre del volcán en texto final
        elif ("volcán" in texto.lower() or "volcan" in texto.lower() or "volcano" in texto.lower()):
            import re
            # Reemplazar TODOS los nombres de volcanes en TODOS los párrafos
            cambio_realizado = False
            for volcan_antiguo in VOLCANES_ACTIVOS:
                if volcan_antiguo.lower() in texto.lower():
                    # Iterar sobre TODOS los párrafos (no solo el primero)
                    for p in shape.text_frame.paragraphs:
                        texto_parrafo = p.text
                        if volcan_antiguo.lower() in texto_parrafo.lower():
                            # Reemplazo case-insensitive
                            patron = re.compile(re.escape(volcan_antiguo), re.IGNORECASE)
                            texto_nuevo_p = patron.sub(volcan_nombre, texto_parrafo)
                            
                            # También reemplazar "mes de [mes] [año]"
                            # Extraer mes y año de fecha_fin
                            from datetime import datetime
                            try:
                                dt = datetime.strptime(fecha_fin, '%Y-%m-%d')
                                mes_actual = MESES_ES[dt.month]
                                ano_actual = dt.year
                                
                                # Buscar patrón "mes de [cualquier mes] [año]"
                                patron_fecha = re.compile(
                                    r'(mes de )\w+(?: \d{4})?',
                                    re.IGNORECASE
                                )
                                texto_nuevo_p = patron_fecha.sub(
                                    f"mes de {mes_actual} {ano_actual}",
                                    texto_nuevo_p
                                )
                            except:
                                pass
                            
                            if texto_nuevo_p != texto_parrafo:
                                # Guardar formato
                                fmt = None
                                if p.runs:
                                    fmt = {'name': p.runs[0].font.name, 'size': p.runs[0].font.size,
                                          'bold': p.runs[0].font.bold, 'italic': p.runs[0].font.italic}
                                
                                # Reemplazar texto
                                p.clear()
                                run = p.add_run()
                                run.text = texto_nuevo_p
                                
                                # Restaurar formato
                                if fmt:
                                    if fmt['name']: run.font.name = fmt['name']
                                    if fmt['size']: run.font.size = fmt['size']
                                    if fmt['bold'] is not None: run.font.bold = fmt['bold']
                                    if fmt['italic'] is not None: run.font.italic = fmt['italic']
                                
                                cambio_realizado = True
            
            if cambio_realizado:
                print(f"       ✅ Nombre volcán actualizado")
                textos_ok += 1
    
    if textos_ok < 2:
        print(f"    ⚠️  Textos: {textos_ok} (esperados 2-3)")
    
    print(f"    🖼️  Reemplazando GIFs...")
    # Encontrar subtitulos para emparejar cada imagen con su tipo correcto.
    # Sortear por top fallaba en la plantilla nueva: las dos imagenes tienen
    # tops casi iguales (2.91 vs 2.92cm) y el sort las invertia.
    subtitle_rgb = None
    subtitle_thermal = None
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        t = s.text_frame.text.lower()
        if "color verdadero" in t and "time lapse" in t:
            subtitle_rgb = s
        elif "falso color" in t and "time lapse" in t:
            subtitle_thermal = s

    shapes_img = [s for s in slide.shapes if s.shape_type == 13]

    def proximidad_centro(pic, sub):
        """Distancia horizontal entre centros de imagen y subtitulo"""
        pic_cx = pic.left + pic.width / 2
        sub_cx = sub.left + sub.width / 2
        return abs(pic_cx - sub_cx)

    if len(shapes_img) >= 2 and subtitle_rgb is not None and subtitle_thermal is not None:
        # Emparejar por proximidad horizontal con subtitulos
        pic_rgb = min(shapes_img, key=lambda p: proximidad_centro(p, subtitle_rgb))
        pic_thermal = next(p for p in shapes_img if p is not pic_rgb)

        rgb_pos = (pic_rgb.left, pic_rgb.top, pic_rgb.width, pic_rgb.height)
        pic_rgb.element.getparent().remove(pic_rgb.element)
        slide.shapes.add_picture(gif_rgb_final, *rgb_pos)
        print(f"       ✅ RGB")

        thermal_pos = (pic_thermal.left, pic_thermal.top, pic_thermal.width, pic_thermal.height)
        pic_thermal.element.getparent().remove(pic_thermal.element)
        slide.shapes.add_picture(gif_thermal_final, *thermal_pos)
        print(f"       ✅ Thermal")
    elif len(shapes_img) >= 2:
        # Fallback: sin subtitulos identificables, usar orden original (top, left)
        print(f"       ⚠️  Subtitulos no identificados, usando fallback por posicion")
        shapes_img_data = sorted(
            [{'shape': s, 'top': s.top, 'left': s.left, 'width': s.width, 'height': s.height} for s in shapes_img],
            key=lambda x: (x['top'], x['left'])
        )
        s = shapes_img_data[0]
        s['shape'].element.getparent().remove(s['shape'].element)
        slide.shapes.add_picture(gif_rgb_final, s['left'], s['top'], s['width'], s['height'])
        s = shapes_img_data[1]
        s['shape'].element.getparent().remove(s['shape'].element)
        slide.shapes.add_picture(gif_thermal_final, s['left'], s['top'], s['width'], s['height'])
    
    carpeta_reportes = os.path.join(OUTPUT_DIR, volcan_nombre, "reportes")
    os.makedirs(carpeta_reportes, exist_ok=True)
    
    output_path = os.path.join(carpeta_reportes, 
                               f"{volcan_nombre}_Evaluacion_Mensual_{fecha_fin[:7]}.pptx")
    
    try:
        prs.save(output_path)
        size_mb = os.path.getsize(output_path) / (1024 * 1024)
        status = "✅" if size_mb < 3.0 else "⚠️"
        print(f"   {status} PPT: {size_mb:.2f} MB")
        
        try:
            os.remove(temp_rgb)
            os.remove(temp_thermal)
        except:
            pass
        
        return output_path
    except Exception as e:
        print(f"    ❌ Error: {e}")
        return None

def generar_ppt_combinado(ppts_generados, fecha_inicio, fecha_fin):
    """Genera un PPT con todos los volcanes (1 slide por volcán).

    Opera directamente sobre los ZIPs (PPTX = ZIP) para evitar problemas de
    referencias cruzadas entre paquetes python-pptx. Para cada PPT individual:
    - Copia el slide XML
    - Copia los archivos de media (imágenes/GIFs) renombrados para evitar colisiones
    - Actualiza el .rels del slide con las nuevas rutas de media
    - Registra el slide en presentation.xml y [Content_Types].xml
    """
    if not ppts_generados:
        print("⚠️  Sin PPTs individuales para combinar")
        return None

    print(f"\n{'='*80}")
    print(f"📚 GENERANDO PPT COMBINADO ({len(ppts_generados)} volcanes)...")

    os.makedirs("docs/reportes", exist_ok=True)
    output_path = f"docs/reportes/Evaluacion_Completa_{fecha_inicio}_{fecha_fin}.pptx"

    # Trabajamos en un directorio temporal que iremos llenando
    tmpdir = tempfile.mkdtemp()

    try:
        # ── PASO 1: Extraer el primer PPTX como base ──────────────────────────
        with zipfile.ZipFile(ppts_generados[0], 'r') as z:
            z.extractall(tmpdir)

        # Leer presentation.xml para saber cuántos slides hay ya
        prs_xml_path = os.path.join(tmpdir, 'ppt', 'presentation.xml')
        prs_rels_path = os.path.join(tmpdir, 'ppt', '_rels', 'presentation.xml.rels')
        ct_path = os.path.join(tmpdir, '[Content_Types].xml')

        NS_PRS  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        NS_R    = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        NS_CT   = 'http://schemas.openxmlformats.org/package/2006/content-types'
        NS_RELS = 'http://schemas.openxmlformats.org/package/2006/relationships'

        # Contar slides existentes en el base
        prs_tree = etree.parse(prs_xml_path)
        sldIdLst = prs_tree.find(f'{{{NS_PRS}}}sldIdLst')
        next_slide_num = len(sldIdLst) + 1          # próximo índice de slide
        next_slide_id  = max(                        # próximo id numérico
            int(el.get('id')) for el in sldIdLst
        ) + 1

        # Contar media existente (para no colisionar nombres)
        media_dir = os.path.join(tmpdir, 'ppt', 'media')
        os.makedirs(media_dir, exist_ok=True)
        next_media_num = len(os.listdir(media_dir)) + 1

        # ── PASO 2: Para cada PPT adicional, copiar slide + media ─────────────
        for ppt_path in ppts_generados[1:]:
            nombre = os.path.basename(ppt_path)
            try:
                with zipfile.ZipFile(ppt_path, 'r') as z:
                    nombres_zip = z.namelist()

                    # Encontrar el primer slide
                    slides_en_zip = sorted(
                        [n for n in nombres_zip if n.startswith('ppt/slides/slide') and n.endswith('.xml')
                         and '_rels' not in n]
                    )
                    if not slides_en_zip:
                        print(f"   ⚠️  Sin slides: {nombre}")
                        continue

                    slide_zip_path = slides_en_zip[0]
                    slide_rels_zip = f"ppt/slides/_rels/{os.path.basename(slide_zip_path)}.rels"

                    # Leer XML del slide y de su .rels
                    slide_xml_bytes = z.read(slide_zip_path)
                    try:
                        slide_rels_bytes = z.read(slide_rels_zip)
                    except KeyError:
                        slide_rels_bytes = None

                    # Copiar media renombrada y construir mapa old_name -> new_name
                    media_map = {}
                    for item in nombres_zip:
                        if item.startswith('ppt/media/'):
                            old_media_name = os.path.basename(item)
                            ext = os.path.splitext(old_media_name)[1]
                            new_media_name = f"media{next_media_num}{ext}"
                            next_media_num += 1
                            media_map[old_media_name] = new_media_name
                            dest = os.path.join(media_dir, new_media_name)
                            with z.open(item) as src, open(dest, 'wb') as dst:
                                dst.write(src.read())

                    # Actualizar .rels del slide con los nuevos nombres de media
                    if slide_rels_bytes and media_map:
                        rels_tree = etree.fromstring(slide_rels_bytes)
                        for rel_el in rels_tree.findall(f'{{{NS_RELS}}}Relationship'):
                            target = rel_el.get('Target', '')
                            old_name = os.path.basename(target)
                            if old_name in media_map:
                                rel_el.set('Target', f'../media/{media_map[old_name]}')
                        slide_rels_bytes = etree.tostring(rels_tree, xml_declaration=True,
                                                          encoding='UTF-8', standalone=True)

                # Escribir slide en el directorio base
                new_slide_name = f"slide{next_slide_num}.xml"
                slide_dest = os.path.join(tmpdir, 'ppt', 'slides', new_slide_name)
                with open(slide_dest, 'wb') as f:
                    f.write(slide_xml_bytes)

                # Escribir .rels del slide
                if slide_rels_bytes:
                    rels_dir = os.path.join(tmpdir, 'ppt', 'slides', '_rels')
                    os.makedirs(rels_dir, exist_ok=True)
                    with open(os.path.join(rels_dir, f"{new_slide_name}.rels"), 'wb') as f:
                        f.write(slide_rels_bytes)

                # Registrar slide en presentation.xml
                new_rId = f"rId_slide{next_slide_num}"
                etree.SubElement(sldIdLst, f'{{{NS_PRS}}}sldId',
                                 id=str(next_slide_id), r_id=new_rId)
                # Fijar el atributo con namespace correcto
                sldIdLst[-1].set(f'{{{NS_R}}}id', new_rId)
                sldIdLst[-1].attrib.pop('r_id', None)

                # Registrar en presentation.xml.rels
                rels_tree_prs = etree.parse(prs_rels_path)
                etree.SubElement(
                    rels_tree_prs.getroot(),
                    f'{{{NS_RELS}}}Relationship',
                    Id=new_rId,
                    Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide',
                    Target=f'slides/{new_slide_name}'
                )
                rels_tree_prs.write(prs_rels_path, xml_declaration=True,
                                    encoding='UTF-8', standalone=True)

                # Registrar slide en [Content_Types].xml
                ct_tree = etree.parse(ct_path)
                etree.SubElement(
                    ct_tree.getroot(),
                    f'{{{NS_CT}}}Override',
                    PartName=f'/ppt/slides/{new_slide_name}',
                    ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml'
                )
                ct_tree.write(ct_path, xml_declaration=True,
                              encoding='UTF-8', standalone=True)

                next_slide_num += 1
                next_slide_id  += 1
                print(f"   ✅ {nombre}")

            except Exception as e:
                print(f"   ❌ Error copiando {nombre}: {e}")
                import traceback; traceback.print_exc()

        # Guardar presentation.xml actualizado
        prs_tree.write(prs_xml_path, xml_declaration=True,
                       encoding='UTF-8', standalone=True)

        # ── PASO 3: Re-empaquetar el directorio como ZIP/PPTX ────────────────
        if os.path.exists(output_path):
            os.remove(output_path)

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for root, dirs, files in os.walk(tmpdir):
                for file in files:
                    full_path = os.path.join(root, file)
                    arcname = os.path.relpath(full_path, tmpdir)
                    zout.write(full_path, arcname)

        size_mb = os.path.getsize(output_path) / (1024 * 1024)
        print(f"✅ PPT combinado: {output_path} ({size_mb:.1f} MB)")
        return output_path

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def main():
    print("="*80)
    print("📊 PPT GENERATOR V5.2")
    print("   NUEVO: PPT combinado con todos los volcanes")
    print("="*80)

    if not os.path.exists(PLANTILLA_PATH):
        print(f"\n❌ Plantilla no encontrada: {PLANTILLA_PATH}")
        return

    volcan_env = os.getenv('VOLCAN')
    fecha_inicio_env = os.getenv('FECHA_INICIO')
    fecha_fin_env = os.getenv('FECHA_FIN')

    if volcan_env:
        print(f"\n📋 MODO MANUAL: Procesando {volcan_env}")
        volcanes_a_procesar = [volcan_env]
    else:
        print("\n🤖 MODO AUTOMÁTICO: Procesando todos los volcanes")
        volcanes_a_procesar = VOLCANES_ACTIVOS

    ppts = []
    for volcan in volcanes_a_procesar:
        try:
            ppt = generar_ppt(volcan)
            if ppt:
                ppts.append(ppt)
        except Exception as e:
            print(f"❌ Error en {volcan}: {e}")

    print("\n" + "="*80)
    print(f"✅ {len(ppts)} PPTs individuales generados")
    for ppt in ppts:
        size_mb = os.path.getsize(ppt) / (1024 * 1024)
        print(f"   📄 {os.path.basename(ppt)}: {size_mb:.2f} MB")

    # Generar PPT combinado (solo si hay más de 1 volcán)
    if len(ppts) > 1:
        # Usar fechas de env vars, o derivar del primer GIF disponible
        if not fecha_inicio_env or not fecha_fin_env:
            # Extraer del nombre del primer PPT generado
            partes = os.path.basename(ppts[0]).replace('.pptx', '').split('_')
            fecha_inicio_env = partes[-2] if len(partes) >= 2 else 'fecha'
            fecha_fin_env = partes[-1] if len(partes) >= 1 else 'fin'
        generar_ppt_combinado(ppts, fecha_inicio_env, fecha_fin_env)

    print("="*80)

if __name__ == "__main__":
    main()
