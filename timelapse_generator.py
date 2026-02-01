"""
PPT_GENERATOR.PY
Genera presentación PowerPoint mensual con GIFs de timelapses
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import os
import glob
from io import BytesIO

# =========================
# CONFIGURACIÓN
# =========================

VOLCANES_ACTIVOS = ["Villarrica", "Llaima"]
TEMPLATE_PPT = "data/Cambios_morfologicos.pptx"  # Plantilla en el repositorio

# =========================
# GENERADOR PPT
# =========================

def generar_ppt_mensual(volcan_nombre, mes=None, año=None):
    """
    Genera PPT mensual para un volcán
    
    Args:
        volcan_nombre: Nombre del volcán
        mes: Mes (1-12), si None usa mes anterior
        año: Año, si None usa año actual
    
    Returns:
        str: Path al PPT generado
    """
    
    # Si no se especifica mes/año, usar mes anterior
    if mes is None or año is None:
        ahora = datetime.now()
        if ahora.month == 1:
            mes = 12
            año = ahora.year - 1
        else:
            mes = ahora.month - 1
            año = ahora.year
    
    print(f"\n📊 Generando PPT: {volcan_nombre} - {año}-{mes:02d}")
    
    # Cargar plantilla
    if not os.path.exists(TEMPLATE_PPT):
        print(f"   ❌ Plantilla no encontrada: {TEMPLATE_PPT}")
        return None
    
    prs = Presentation(TEMPLATE_PPT)
    
    # Buscar GIFs del mes
    carpeta_timelapses = f"data/sentinel2/{volcan_nombre}/timelapses"
    mes_str = f"{año}-{mes:02d}"
    
    gif_rgb_path = f"{carpeta_timelapses}/{volcan_nombre}_RGB_{mes_str}.gif"
    gif_thermal_path = f"{carpeta_timelapses}/{volcan_nombre}_ThermalFalseColor_{mes_str}.gif"
    
    if not os.path.exists(gif_rgb_path) or not os.path.exists(gif_thermal_path):
        print(f"   ⚠️ No se encontraron GIFs para {mes_str}")
        print(f"      Buscado: {gif_rgb_path}")
        return None
    
    print(f"   ✅ GIF RGB: {gif_rgb_path}")
    print(f"   ✅ GIF Thermal: {gif_thermal_path}")
    
    # Obtener rango de fechas de los GIFs
    carpeta_rgb = f"data/sentinel2/{volcan_nombre}/RGB"
    imagenes_mes = sorted(glob.glob(f"{carpeta_rgb}/{año}-{mes:02d}-*.png"))
    
    if imagenes_mes:
        fecha_inicio = os.path.basename(imagenes_mes[0]).split('_')[0]
        fecha_fin = os.path.basename(imagenes_mes[-1]).split('_')[0]
        
        # Formato: "02 diciembre – 30 diciembre 2025"
        meses_es = {
            1: 'enero', 2: 'febrero', 3: 'marzo', 4: 'abril',
            5: 'mayo', 6: 'junio', 7: 'julio', 8: 'agosto',
            9: 'septiembre', 10: 'octubre', 11: 'noviembre', 12: 'diciembre'
        }
        
        dia_inicio = int(fecha_inicio.split('-')[2])
        dia_fin = int(fecha_fin.split('-')[2])
        mes_nombre = meses_es[mes]
        
        rango_fechas_rgb = f"Imágenes Sentinel 2 L2A color verdadero\nTime Lapse {dia_inicio:02d} {mes_nombre} – {dia_fin:02d} {mes_nombre} {año}"
        rango_fechas_thermal = f"Imágenes Sentinel 2 L2A falso color térmico\nTime Lapse {dia_inicio:02d} {mes_nombre} – {dia_fin:02d} {mes_nombre} {año}"
    else:
        rango_fechas_rgb = f"Imágenes Sentinel 2 L2A color verdadero\nTime Lapse {mes_str}"
        rango_fechas_thermal = f"Imágenes Sentinel 2 L2A falso color térmico\nTime Lapse {mes_str}"
    
    print(f"   📅 Rango: {rango_fechas_rgb.split('Time Lapse ')[1]}")
    
    # Modificar slide
    slide = prs.slides[0]
    
    # SHAPE 1: Título - "Cambios Morfológicos"
    # SHAPE 2: Imagen RGB (izquierda) - posición x=0.49"
    # SHAPE 3: Imagen Thermal (derecha) - posición x=6.69"
    # SHAPE 4: Texto evaluación (abajo)
    # SHAPE 5: Texto RGB (título izquierda) - posición x=5.15"
    # SHAPE 6: Texto Thermal (título izquierda) - posición x=-1.42"
    
    shapes_to_remove = []
    
    for idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "text"):
            # Título principal
            if "Cambios Morfológicos" in shape.text:
                shape.text = f"Cambios Morfológicos - {volcan_nombre}"
                print(f"   ✅ Actualizado título")
            
            # Identificar textos por posición
            elif hasattr(shape, 'left'):
                x_pos = shape.left.inches
                
                # Texto derecho (x > 5) = RGB
                if x_pos > 5:
                    shape.text = rango_fechas_rgb
                    print(f"   ✅ Actualizado texto RGB (x={x_pos:.2f}\")")
                
                # Texto izquierdo (x < 5) = Thermal
                elif "Time Lapse" in shape.text or "Imágenes Sentinel" in shape.text:
                    shape.text = rango_fechas_thermal
                    print(f"   ✅ Actualizado texto Thermal (x={x_pos:.2f}\")")
        
        # Marcar imágenes para reemplazo
        if shape.shape_type == 13:  # Picture
            shapes_to_remove.append((idx, shape))
    
    # Reemplazar imágenes (de atrás hacia adelante para no afectar índices)
    for idx, shape in reversed(shapes_to_remove):
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        x_pos = left.inches
        
        # Imagen izquierda (x < 4) = RGB
        # Imagen derecha (x > 4) = Thermal
        if x_pos < 4:
            gif_path = gif_rgb_path
            print(f"   🖼️ Reemplazando imagen izquierda (x={x_pos:.2f}\") con RGB")
        else:
            gif_path = gif_thermal_path
            print(f"   🖼️ Reemplazando imagen derecha (x={x_pos:.2f}\") con Thermal")
        
        # Eliminar imagen antigua
        sp = shape.element
        sp.getparent().remove(sp)
        
        # Agregar GIF nuevo en la misma posición
        slide.shapes.add_picture(gif_path, left, top, width, height)
    
    # Guardar PPT
    carpeta_output = f"data/sentinel2/{volcan_nombre}/reportes"
    os.makedirs(carpeta_output, exist_ok=True)
    
    output_path = f"{carpeta_output}/{volcan_nombre}_Evaluacion_Mensual_{año}-{mes:02d}.pptx"
    prs.save(output_path)
    
    size_mb = os.path.getsize(output_path) / (1024 * 1024)
    print(f"   ✅ PPT generado: {output_path}")
    print(f"   📦 Tamaño: {size_mb:.2f} MB")
    
    return output_path


# =========================
# PROCESO PRINCIPAL
# =========================

def main():
    print("="*80)
    print("📊 GENERADOR PPT EVALUACIÓN MENSUAL")
    print("="*80)
    
    ppts_generados = []
    
    for volcan in VOLCANES_ACTIVOS:
        ppt_path = generar_ppt_mensual(volcan)
        if ppt_path:
            ppts_generados.append(ppt_path)
    
    print("\n" + "="*80)
    print(f"✅ PROCESO COMPLETADO - {len(ppts_generados)} PPTs generados")
    print("="*80)
    
    for ppt in ppts_generados:
        print(f"   📁 {ppt}")


if __name__ == "__main__":
    main()"""
TIMELAPSE_GENERATOR.PY
Genera GIF animado con todas las imágenes del último mes
"""

import os
from PIL import Image
from datetime import datetime
import glob

# =========================
# CONFIGURACIÓN
# =========================

VOLCANES_ACTIVOS = ["Villarrica", "Llaima"]  # Expandir cuando se activen más

FPS = 1  # Frame por segundo
DURACION_FRAME = 1000  # Milisegundos por frame (1000ms = 1s)

# =========================
# GENERADOR DE GIF
# =========================

def generar_gif(volcan_nombre, tipo='RGB'):
    """
    Genera GIF timelapse para un volcán y tipo de imagen
    
    Args:
        volcan_nombre: Nombre del volcán
        tipo: 'RGB' o 'ThermalFalseColor'
    
    Returns:
        str: Path al GIF generado
    """
    
    print(f"\n🎬 Generando GIF: {volcan_nombre} - {tipo}")
    
    # Carpeta de imágenes
    carpeta_imagenes = f"data/sentinel2/{volcan_nombre}/{tipo}"
    
    if not os.path.exists(carpeta_imagenes):
        print(f"   ❌ Carpeta no existe: {carpeta_imagenes}")
        return None
    
    # Obtener todas las imágenes PNG
    imagenes_paths = sorted(glob.glob(f"{carpeta_imagenes}/*.png"))
    
    if len(imagenes_paths) == 0:
        print(f"   ⚠️ No hay imágenes en {carpeta_imagenes}")
        return None
    
    print(f"   📷 Encontradas {len(imagenes_paths)} imágenes")
    
    # Cargar imágenes
    imagenes = []
    fechas = []
    
    for img_path in imagenes_paths:
        try:
            img = Image.open(img_path)
            
            # Extraer fecha del nombre del archivo
            # Formato: YYYY-MM-DD_RGB.png
            nombre_archivo = os.path.basename(img_path)
            fecha = nombre_archivo.split('_')[0]
            fechas.append(fecha)
            
            # Agregar texto de fecha a la imagen
            img_con_texto = agregar_fecha_a_imagen(img, fecha, tipo)
            imagenes.append(img_con_texto)
            
            print(f"   ✅ {fecha}")
        except Exception as e:
            print(f"   ❌ Error cargando {img_path}: {e}")
            continue
    
    if len(imagenes) == 0:
        print(f"   ❌ No se pudieron cargar imágenes")
        return None
    
    # Crear carpeta de salida
    carpeta_gif = f"data/sentinel2/{volcan_nombre}/timelapses"
    os.makedirs(carpeta_gif, exist_ok=True)
    
    # Generar GIF
    fecha_inicio = fechas[0]
    fecha_fin = fechas[-1]
    
    # Usar mes del período más reciente para nombre consistente
    mes_actual = fecha_fin[:7]  # YYYY-MM
    output_path = f"{carpeta_gif}/{volcan_nombre}_{tipo}_{mes_actual}.gif"
    
    try:
        imagenes[0].save(
            output_path,
            save_all=True,
            append_images=imagenes[1:],
            duration=DURACION_FRAME,
            loop=0,  # Loop infinito
            optimize=True
        )
        
        size_mb = os.path.getsize(output_path) / (1024 * 1024)
        print(f"   ✅ GIF generado: {output_path}")
        print(f"   📦 Tamaño: {size_mb:.2f} MB")
        print(f"   🎞️ Frames: {len(imagenes)}")
        
        return output_path
    
    except Exception as e:
        print(f"   ❌ Error generando GIF: {e}")
        return None


def agregar_fecha_a_imagen(img, fecha, tipo):
    """
    Agrega texto de fecha y tipo a la imagen
    
    Args:
        img: PIL Image
        fecha: String de fecha (YYYY-MM-DD)
        tipo: 'RGB' o 'ThermalFalseColor'
    
    Returns:
        PIL Image con texto
    """
    from PIL import ImageDraw, ImageFont
    
    # Crear copia
    img_copy = img.copy()
    draw = ImageDraw.Draw(img_copy)
    
    # Intentar usar fuente del sistema, si no, usar default
    try:
        font_fecha = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
        font_tipo = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
    except:
        font_fecha = ImageFont.load_default()
        font_tipo = ImageFont.load_default()
    
    # Texto de fecha
    texto_fecha = f"📅 {fecha}"
    
    # Texto de tipo
    tipo_texto = "RGB (Color Real)" if tipo == 'RGB' else "Falso Color Térmico (SWIR)"
    
    # Posición del texto (esquina superior izquierda)
    x = 20
    y = 20
    
    # Dibujar fondo semi-transparente para el texto
    # (PIL no soporta transparencia directamente, usamos un rectángulo negro)
    bbox_fecha = draw.textbbox((x, y), texto_fecha, font=font_fecha)
    draw.rectangle(bbox_fecha, fill=(0, 0, 0))
    
    bbox_tipo = draw.textbbox((x, y + 35), tipo_texto, font=font_tipo)
    draw.rectangle(bbox_tipo, fill=(0, 0, 0))
    
    # Dibujar texto
    draw.text((x, y), texto_fecha, fill=(255, 255, 255), font=font_fecha)
    draw.text((x, y + 35), tipo_texto, fill=(200, 200, 200), font=font_tipo)
    
    return img_copy


# =========================
# PROCESO PRINCIPAL
# =========================

def main():
    print("="*80)
    print("🎬 GENERADOR DE TIMELAPSE GIF")
    print("="*80)
    
    gifs_generados = []
    
    for volcan in VOLCANES_ACTIVOS:
        print(f"\n🌋 Procesando: {volcan}")
        
        # Generar GIF RGB
        gif_rgb = generar_gif(volcan, 'RGB')
        if gif_rgb:
            gifs_generados.append(gif_rgb)
        
        # Generar GIF Thermal
        gif_thermal = generar_gif(volcan, 'ThermalFalseColor')
        if gif_thermal:
            gifs_generados.append(gif_thermal)
    
    print("\n" + "="*80)
    print(f"✅ PROCESO COMPLETADO - {len(gifs_generados)} GIFs generados")
    print("="*80)
    
    for gif in gifs_generados:
        print(f"   📁 {gif}")


if __name__ == "__main__":
    main()
